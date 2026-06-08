"""
services/ingest.py

Document ingestion business logic.
Handles: reading files, chunking, embedding, storing in vector DB.
No FastAPI/Streamlit imports — pure Python, fully testable.
"""

import asyncio
from pathlib import Path
from typing import Generator

from langchain_text_splitters import RecursiveCharacterTextSplitter
from core.config import settings


# ── Text splitter (shared) ─────────────────────────────────────────────────────

def get_splitter() -> RecursiveCharacterTextSplitter:
    return RecursiveCharacterTextSplitter(
        chunk_size=settings.chunk_size,
        chunk_overlap=settings.chunk_overlap,
        separators=["\n\n", "\n", ". ", " ", ""],
    )


# ── File readers ───────────────────────────────────────────────────────────────

def read_txt(path: Path) -> str:
    return path.read_text(errors="ignore")


def read_pdf(path: Path) -> str:
    import pypdf

    text = "\n".join(p.extract_text() or "" for p in pypdf.PdfReader(str(path)).pages)
    if text.strip():
        return text
    # No text layer (scanned / image-only PDF) → fall back to OCR.
    return _ocr_pdf(path)


def _ocr_pdf(path: Path) -> str:
    """OCR a scanned PDF. Returns '' if OCR tooling isn't available."""
    try:
        import pytesseract
        from pdf2image import convert_from_path
    except ImportError:
        return ""
    try:
        pages = convert_from_path(str(path), dpi=200)
        return "\n".join(pytesseract.image_to_string(img) for img in pages)
    except Exception:
        # Missing system binaries (tesseract/poppler) or unreadable scan.
        return ""


def _ocr_embedded_images(path: Path, media_prefix: str) -> str:
    """
    OCR images embedded inside an Office Open XML file (docx/pptx are zips).
    Returns '' if OCR tooling isn't available or there are no images.
    """
    import io
    import zipfile

    try:
        import pytesseract
        from PIL import Image
    except ImportError:
        return ""

    exts = (".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif", ".gif")
    out: list[str] = []
    try:
        with zipfile.ZipFile(path) as z:
            for name in z.namelist():
                if not name.startswith(media_prefix) or not name.lower().endswith(exts):
                    continue
                try:
                    text = pytesseract.image_to_string(Image.open(io.BytesIO(z.read(name))))
                    if text.strip():
                        out.append(text)
                except Exception:
                    continue
    except Exception:
        return ""
    return "\n".join(out)


def read_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    parts: list[str] = [p.text for p in doc.paragraphs if p.text.strip()]

    # Table cell text (commonly the bulk of real-world docx content).
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))

    # Headers / footers.
    for section in doc.sections:
        for container in (section.header, section.footer):
            for p in container.paragraphs:
                if p.text.strip():
                    parts.append(p.text)

    text = "\n".join(parts)
    # OCR any embedded images (scans pasted into the doc).
    ocr = _ocr_embedded_images(path, "word/media/")
    return f"{text}\n{ocr}".strip() if ocr else text


def read_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []

    def walk(shapes):
        for shape in shapes:
            if shape.shape_type == 6 and hasattr(shape, "shapes"):  # grouped shapes
                walk(shape.shapes)
                continue
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))

    for slide in prs.slides:
        walk(slide.shapes)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                parts.append(notes)

    text = "\n".join(parts)
    ocr = _ocr_embedded_images(path, "ppt/media/")
    return f"{text}\n{ocr}".strip() if ocr else text


# Legacy binary formats python-docx / python-pptx can't read — reject clearly.
_LEGACY = {".doc": "Word (.doc)", ".ppt": "PowerPoint (.ppt)"}


def read_file(path: Path) -> str:
    """Read any supported file type."""
    suffix = path.suffix.lower()
    if suffix in _LEGACY:
        modern = ".docx" if suffix == ".doc" else ".pptx"
        raise ValueError(
            f"Legacy {_LEGACY[suffix]} format not supported — save as {modern} and re-upload."
        )
    readers = {
        ".txt":  read_txt,
        ".md":   read_txt,
        ".pdf":  read_pdf,
        ".docx": read_docx,
        ".pptx": read_pptx,
    }
    reader = readers.get(suffix)
    if reader is None:
        raise ValueError(f"Unsupported file type: {path.suffix}")
    return reader(path)


def read_folder(folder: str) -> Generator[tuple[str, str], None, None]:
    """Yield (filename, text) for every supported file in a folder."""
    for path in sorted(Path(folder).rglob("*")):
        if path.suffix.lower() not in {".txt", ".md", ".pdf", ".docx", ".pptx"}:
            continue
        text = read_file(path)
        if text and text.strip():
            yield path.name, text


def read_google_drive(folder_id: str) -> Generator[tuple[str, str], None, None]:
    """Yield (filename, text) for files in a Google Drive folder."""
    import io
    from google.oauth2 import service_account
    from googleapiclient.discovery import build
    from googleapiclient.http import MediaIoBaseDownload
    from core.config import settings
    import os

    sa_file = os.environ["GOOGLE_SERVICE_ACCOUNT_JSON"]
    creds   = service_account.Credentials.from_service_account_file(
        sa_file, scopes=["https://www.googleapis.com/auth/drive.readonly"]
    )
    drive   = build("drive", "v3", credentials=creds)
    results = drive.files().list(
        q=f"'{folder_id}' in parents and trashed=false",
        fields="files(id,name,mimeType)"
    ).execute()

    for f in results.get("files", []):
        mime, fid, name = f["mimeType"], f["id"], f["name"]
        try:
            if mime == "application/vnd.google-apps.document":
                data = drive.files().export(fileId=fid, mimeType="text/plain").execute()
                yield name, data.decode("utf-8", errors="ignore")
            elif mime in ("text/plain", "text/markdown"):
                buf = io.BytesIO()
                MediaIoBaseDownload(buf, drive.files().get_media(fileId=fid)).next_chunk()
                yield name, buf.getvalue().decode("utf-8", errors="ignore")
        except Exception as exc:
            print(f"  Error reading {name}: {exc}")


# ── Core ingestion logic ───────────────────────────────────────────────────────

def ingest_texts(
    source_iter: Generator[tuple[str, str], None, None],
    on_progress: callable = None,
) -> dict:
    """
    Ingest (filename, text) pairs into the vector store.

    Args:
        source_iter:  Generator yielding (filename, text) tuples
        on_progress:  Optional callback(filename, chunks_count, done, total)

    Returns:
        {"files": int, "chunks": int}
    """
    from providers.embeddings import get_embedder
    from providers.vectorstore import get_vectorstore

    splitter = get_splitter()
    embedder = get_embedder()
    store    = get_vectorstore()

    total_files  = 0
    total_chunks = 0

    for filename, text in source_iter:
        chunks = splitter.split_text(text)
        if not chunks:
            continue

        embeddings = embedder.embed_documents(chunks)
        store.upsert(chunks, embeddings, source=filename)

        total_files  += 1
        total_chunks += len(chunks)

        if on_progress:
            on_progress(filename, len(chunks))

    return {"files": total_files, "chunks": total_chunks}


async def ingest_folder_async(
    folder: str,
    job: dict,
    log: callable,
) -> None:
    """
    Async ingestion for a folder — used by the API background task.
    Updates the job dict in-place as progress is made.
    """
    from providers.embeddings import get_embedder
    from providers.vectorstore import get_vectorstore

    splitter = get_splitter()

    job["status"] = "running"
    log("Job started")

    try:
        embedder = get_embedder()
        store    = get_vectorstore()
        log("Providers ready ✓")

        files = list(Path(folder).iterdir())
        job["files_found"] = len(files)
        log(f"Processing {len(files)} file(s) …")

        for i, path in enumerate(files, 1):
            fname = path.name
            job["current_file"] = fname
            log(f"[{i}/{len(files)}] {fname}")

            try:
                text = read_file(path)
            except ValueError as exc:
                reason = str(exc) or "unsupported file type"
                log(f"  ↳ skipped ({reason})")
                job.setdefault("skipped", []).append({"file": fname, "reason": reason})
                job["files_done"] += 1
                continue

            chunks = splitter.split_text(text)
            if not chunks:
                log("  ↳ no readable text found, skipped (scanned PDF without OCR?)")
                job.setdefault("skipped", []).append({"file": fname, "reason": "no readable text found"})
                job["files_done"] += 1
                continue

            log(f"  ↳ {len(chunks)} chunks — embedding …")
            embeddings = await asyncio.get_event_loop().run_in_executor(
                None, embedder.embed_documents, chunks
            )

            store.upsert(chunks, embeddings, source=fname)
            job["chunks_total"]  += len(chunks)
            job["files_stored"]   = job.get("files_stored", 0) + 1
            job["files_done"]    += 1
            log(f"  ↳ stored {len(chunks)} chunks ✓")

        job["current_file"] = None
        job["status"]       = "done"
        log(f"All done — {job['chunks_total']:,} chunks stored ✓")

    except Exception as exc:
        job["status"] = "failed"
        job["error"]  = str(exc)
        log(f"ERROR: {exc}")
